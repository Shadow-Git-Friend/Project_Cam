#!/usr/bin/env python3
"""Pass 5: fix caption overlap on slide 7; tidy slide 10 table caption position."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path("/home/hanush/Desktop/Project_Cam")
SRC = ROOT / "presentation_defense_improved_github_prioritized.pptx"

GRAY = RGBColor(0x8A, 0x8A, 0x8A)


def replace_caption(slide, substr, left, top, width, size=Pt(10)):
    # Remove old
    to_rm = [sh for sh in slide.shapes
             if sh.has_text_frame and substr in sh.text_frame.text
             and not sh.has_table]
    for s in to_rm:
        sp = s._element
        sp.getparent().remove(sp)
    # Add new at correct position
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.22))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    # The substring is enough to identify which caption to restore
    text_map = {
        "Table 2.": "Table 2. Indicative bill of materials (perception + actuation).",
        "Table 3.": "Table 3. Evaluation matrix: experiment type, input, metric, and outcome.",
    }
    r.text = text_map.get(substr, substr)
    r.font.size = size
    r.font.italic = True
    r.font.color.rgb = GRAY


def main():
    prs = Presentation(str(SRC))

    # Slide 7: place Table 2 caption just above BOM header (which starts at T=4.93)
    s7 = prs.slides[6]
    replace_caption(s7, "Table 2.", Inches(7.50), Inches(4.68), Inches(5.28), size=Pt(10))

    # Slide 10: reposition Table 3 caption above the trimmed 5-row matrix; header row is at T≈2.29
    # caption should be at T≈2.03 (already there) but some row-rectangles were left behind.
    # Do a second sweep to remove any orphan row rectangles below the 4-row data (y > 4.95).
    s10 = prs.slides[9]
    # After trimming, keep everything at y<=5.15. Remove rectangles beyond that.
    to_rm = []
    for shape in s10.shapes:
        if shape.shape_type == 1:  # AUTO_SHAPE
            # rectangle in the grid area?
            top_in = shape.top / 914400
            h_in = shape.height / 914400
            w_in = shape.width / 914400
            if top_in > 4.95 and top_in < 6.80 and h_in < 0.80 and w_in > 2.5 and w_in < 3.5:
                to_rm.append(shape)
    for s in to_rm:
        sp = s._element
        sp.getparent().remove(sp)
    if to_rm:
        print(f"  removed {len(to_rm)} orphan row rectangles from slide 10")

    prs.save(str(SRC))
    print(f"Saved: {SRC}")


if __name__ == "__main__":
    main()
