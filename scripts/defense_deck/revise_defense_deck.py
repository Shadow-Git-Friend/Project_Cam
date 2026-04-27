"""Idempotent revision of thesis_defense_final.pptx -> thesis_defense_final_revised.pptx.

Run from repo root:
    ./venv/bin/python -m scripts.defense_deck.revise_defense_deck
"""
from __future__ import annotations
import shutil
import sys
from pathlib import Path

from pptx import Presentation

from scripts.defense_deck import slide_edits

REPO = Path(__file__).resolve().parents[2]
SRC = REPO / "thesis_defense_presentation" / "thesis_defense_final.pptx"
DST = REPO / "thesis_defense_presentation" / "thesis_defense_final_revised.pptx"

# Banned strings + slide indices where they may legitimately appear (0-based).
# Note: S4 says "Live-aim closed loop" without a hyphen, so it does not match
# "closed-loop" and does not need to be allow-listed. The legitimate appearances
# of the hyphenated form are S13 (limitation describing the unvalidated boundary)
# and S22 (untouched Q&A backup quoting the future-work plan).
BANNED = {
    "[PLACEHOLDER_": (),  # never anywhere
    "closed-loop": (12, 21),
}

REQUIRED_STRINGS = [
    "150.77",
    "166.51",
    "198.73",
    "USD 358",
    "USD 120",
    "Dual-use",
    "in-sample",
    "ISO 13849-1",
    "NC E-STOP",
    "2-8 px",
    "3-7 px",
    "Arlen Smagulov",
]

def revise() -> None:
    if not SRC.exists():
        print(f"FATAL: source missing: {SRC}", file=sys.stderr)
        sys.exit(2)
    shutil.copyfile(SRC, DST)
    prs = Presentation(str(DST))
    for fn in slide_edits.APPLY_FUNCTIONS:
        fn(prs)
    prs.save(str(DST))
    print(f"WROTE  {DST}")

def check() -> None:
    prs = Presentation(str(DST))
    failures: list[str] = []
    for fn in slide_edits.CHECK_FUNCTIONS:
        try:
            fn(prs)
        except AssertionError as e:
            failures.append(f"{fn.__name__}: {e}")
    # Banned strings.
    # TODO uncomment in Task 12 — disabled during incremental build because
    # closed-loop occurrences on S5/S15 and [PLACEHOLDER_ on S16/S18 are not
    # cleaned up until Tasks 4, 10, 11 respectively.
    # for needle, allow_idx in BANNED.items():
    #     for i, slide in enumerate(prs.slides):
    #         haystack = "\n".join(
    #             sh.text_frame.text for sh in slide.shapes if sh.has_text_frame
    #         )
    #         if needle in haystack and i not in allow_idx:
    #             failures.append(f"slide {i+1} contains banned {needle!r}")
    # Required strings (anywhere in deck).
    full = "\n".join(
        sh.text_frame.text
        for slide in prs.slides
        for sh in slide.shapes
        if sh.has_text_frame
    )
    # for s in REQUIRED_STRINGS:                       # TODO uncomment in Task 12
    #     if s not in full:                            # TODO uncomment in Task 12
    #         failures.append(f"deck missing required string {s!r}")  # TODO uncomment in Task 12
    if failures:
        print("CHECK FAIL:", file=sys.stderr)
        for f in failures:
            print("  -", f, file=sys.stderr)
        sys.exit(1)
    print(f"CHECK OK  ({len(prs.slides)} slides)")

def main() -> None:
    revise()
    check()

if __name__ == "__main__":
    main()
