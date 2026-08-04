from typing import Optional
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.util import Pt

def find_shape_by_name(slide: Slide, name: str) -> BaseShape:
    """Return first shape on slide whose .name matches exactly. Raises KeyError if none."""
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    raise KeyError(f"shape {name!r} not on slide {slide.slide_id}; have: {[s.name for s in slide.shapes]}")

def find_shape_containing(slide: Slide, substring: str) -> BaseShape:
    """Return first shape with text_frame whose text contains substring. Raises KeyError."""
    for sh in slide.shapes:
        if sh.has_text_frame and substring in sh.text_frame.text:
            return sh
    raise KeyError(f"no shape on slide {slide.slide_id} contains {substring!r}")

def shape_text(sh: BaseShape) -> str:
    return sh.text_frame.text if sh.has_text_frame else ""

def replace_run_text(sh: BaseShape, search: str, replace: str) -> int:
    """Replace `search` with `replace` inside any single run of any paragraph of sh.
    Preserves run formatting. Returns count of replacements made.
    Idempotent: re-running with same args after a successful replace returns 0."""
    if not sh.has_text_frame:
        return 0
    n = 0
    for para in sh.text_frame.paragraphs:
        for run in para.runs:
            if search in run.text:
                run.text = run.text.replace(search, replace)
                n += 1
    return n

def set_paragraphs(sh: BaseShape, lines: list[str]) -> None:
    """Replace text frame body with `lines` (one paragraph per line). Preserves
    the first paragraph's run-zero formatting and applies it to all new lines.
    Use only when the existing paragraph structure can be discarded."""
    if not sh.has_text_frame:
        raise ValueError(f"shape {sh.name} has no text frame")
    tf = sh.text_frame
    template_para = tf.paragraphs[0]
    template_run = template_para.runs[0] if template_para.runs else None
    template_font = template_run.font if template_run else None
    # Clear all but the first paragraph; clear first paragraph's runs.
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    first = tf.paragraphs[0]
    for r in list(first.runs):
        r._r.getparent().remove(r._r)
    # Write lines.
    first.text = lines[0]
    if template_font is not None and first.runs:
        _copy_font(template_font, first.runs[0].font)
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        if template_font is not None and p.runs:
            _copy_font(template_font, p.runs[0].font)

def _copy_font(src, dst) -> None:
    if src.size is not None:
        dst.size = src.size
    if src.bold is not None:
        dst.bold = src.bold
    if src.italic is not None:
        dst.italic = src.italic
    if src.name is not None:
        dst.name = src.name
    try:
        if src.color and src.color.rgb is not None:
            dst.color.rgb = src.color.rgb
    except Exception:
        pass

def get_table(sh: BaseShape):
    if not sh.has_table:
        raise ValueError(f"shape {sh.name} has no table")
    return sh.table

def set_cell_text(cell, new_text: str) -> None:
    """Replace a table cell's text while keeping its first run's formatting."""
    tf = cell.text_frame
    template_run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else None
    template_font = template_run.font if template_run else None
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    first = tf.paragraphs[0]
    for r in list(first.runs):
        r._r.getparent().remove(r._r)
    first.text = new_text
    if template_font is not None and first.runs:
        _copy_font(template_font, first.runs[0].font)

def find_table_on_slide(slide: Slide):
    """Return first GraphicFrame containing a table on the slide. Raises KeyError."""
    for sh in slide.shapes:
        if sh.has_table:
            return sh.table
    raise KeyError(f"no table on slide {slide.slide_id}")

def assert_slide_text_contains(slide: Slide, needle: str) -> None:
    haystack = "\n".join(shape_text(sh) for sh in slide.shapes)
    if needle not in haystack:
        raise AssertionError(f"slide {slide.slide_id} missing {needle!r}; have:\n{haystack[:1200]}")

def assert_slide_text_not_contains(slide: Slide, needle: str) -> None:
    haystack = "\n".join(shape_text(sh) for sh in slide.shapes)
    if needle in haystack:
        raise AssertionError(f"slide {slide.slide_id} still contains banned {needle!r}")
