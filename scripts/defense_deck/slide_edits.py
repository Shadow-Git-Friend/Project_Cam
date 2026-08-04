"""Per-slide apply/check functions. New tasks append to this file and
extend APPLY_FUNCTIONS / CHECK_FUNCTIONS at the bottom."""
from pptx.presentation import Presentation as _Prs
from pptx.util import Pt

from scripts.defense_deck.helpers import (
    find_shape_by_name,
    find_shape_containing,
    replace_run_text,
    set_paragraphs,
    set_cell_text,
    find_table_on_slide,
    shape_text,
    assert_slide_text_contains,
    assert_slide_text_not_contains,
)

def set_text_size(sh, points: int) -> None:
    if not sh.has_text_frame:
        return
    for para in sh.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(points)

# ---------- Slide 1: title — footer name only ----------

def apply_s1(prs: _Prs) -> None:
    slide = prs.slides[0]
    title = find_shape_by_name(slide, "Title 1")
    set_paragraphs(title, [
        "Pose Guided Predictive Ballistics",
        "for Body Part-Targeted Football Training",
    ])
    footer = find_shape_by_name(slide, "FooterText")
    replace_run_text(footer, "Hanush", "Arlen Smagulov")

def check_s1(prs: _Prs) -> None:
    slide = prs.slides[0]
    assert_slide_text_contains(slide, "Body Part-Targeted Football Training")
    footer = find_shape_by_name(slide, "FooterText")
    txt = shape_text(footer)
    if "Arlen Smagulov" not in txt:
        raise AssertionError(f"S1 footer not updated: {txt!r}")
    if "Hanush" in txt:
        raise AssertionError(f"S1 footer still contains 'Hanush': {txt!r}")

# ---------- Slide 2: motivation — add inline references ----------

S2_NEW_LINES = [
    "Fixed launchers cannot react [Lobster Elite, Sec 1.1]",
    "MoCap is costly and marker-based [OptiTrack/Vicon, Sec 2.1]",
    "Training needs joint-specific delivery",
    "Low-cost cameras can bridge the gap",
    "This thesis builds pose-guided aiming",
]

def apply_s2(prs: _Prs) -> None:
    slide = prs.slides[1]
    block = find_shape_by_name(slide, "Rectangle 12")
    set_paragraphs(block, S2_NEW_LINES)

def check_s2(prs: _Prs) -> None:
    slide = prs.slides[1]
    assert_slide_text_contains(slide, "[Lobster Elite, Sec 1.1]")
    assert_slide_text_contains(slide, "[OptiTrack/Vicon, Sec 2.1]")

# ---------- Slide 3: problem & objectives — soften takeaway ----------

S3_NEW_TAKEAWAY = (
    "Takeaway: RQ1/RQ2 met; RQ3/RQ4 are validated only for static/live-aim. "
    "Moving-target firing is future work (Sec 1.3, 6.4)."
)

def apply_s3(prs: _Prs) -> None:
    slide = prs.slides[2]
    takeaway = find_shape_by_name(slide, "t119")
    set_paragraphs(takeaway, [S3_NEW_TAKEAWAY])

def check_s3(prs: _Prs) -> None:
    slide = prs.slides[2]
    assert_slide_text_contains(slide, "Moving-target firing is future work")
    assert_slide_text_not_contains(slide, "all four objectives are met")

# ---------- Slide 4: background table — rename header + rewrite cell ----------

def apply_s4(prs: _Prs) -> None:
    slide = prs.slides[3]
    set_paragraphs(find_shape_by_name(slide, "t101"), [
        "Gap: low-cost live 3D joint perception is rarely linked to a safety-gated launcher (Sec 2.7).",
    ])
    set_paragraphs(find_shape_by_name(slide, "t152"), [""])
    header = find_shape_by_name(slide, "t111")
    set_paragraphs(header, ["Validation scope"])
    set_paragraphs(find_shape_by_name(slide, "t133"), ["Research prototypes [24,25]"])
    set_paragraphs(find_shape_by_name(slide, "t141"), ["Fixed zones; ball-only"])
    set_paragraphs(find_shape_by_name(slide, "t145"), ["≈USD 200 perception"])
    cell = find_shape_by_name(slide, "t151")
    set_paragraphs(cell, ["Aim + controlled single-shot"])

def check_s4(prs: _Prs) -> None:
    slide = prs.slides[3]
    assert_slide_text_contains(slide, "Validation scope")
    assert_slide_text_contains(slide, "Research prototypes [24,25]")
    assert_slide_text_contains(slide, "≈USD 200 perception")
    assert_slide_text_contains(slide, "Aim + controlled single-shot")
    header = find_shape_by_name(slide, "t111")
    if "Closed-loop" in shape_text(header) or "closed loop" in shape_text(header).lower():
        raise AssertionError("S4 column header still 'Closed-loop'")

# ---------- Slide 5: six contributions ----------

S5_CARDS = {
    "card1": "01 Pose-reactive aiming\nLive joints drive BLM aim",
    "card2": "02 4-camera targeting\n95.17 mm mean (Table 5.1)",
    "card3": "03 Real-time pose\n6.2 ms/batch TRT",
    "card4": "04 Safety validation\nS0-S4; E-STOP <100 ms (Sec 3.14.3)",
    "card5": "05 Voice interface\nASR via UDP gates",
    "card6": "06 Datasets + logs\n36 ball pts + 81 joint trials",
}

def apply_s5(prs: _Prs) -> None:
    slide = prs.slides[4]
    for name, text in S5_CARDS.items():
        sh = find_shape_by_name(slide, name)
        set_paragraphs(sh, text.split("\n"))
        set_text_size(sh, 15)

def check_s5(prs: _Prs) -> None:
    slide = prs.slides[4]
    card1 = find_shape_by_name(slide, "card1")
    if "closed-loop" in shape_text(card1).lower():
        raise AssertionError(f"S5 card1 still contains 'closed-loop': {shape_text(card1)!r}")
    assert_slide_text_contains(slide, "Live joints drive BLM aim")
    assert_slide_text_contains(slide, "Table 5.1")
    assert_slide_text_contains(slide, "Sec 3.14.3")

# ---------- Slide 7: runtime wording — avoid overclaiming end-to-end latency ----------

S7_HEADLINE = (
    "Runtime target: 15 FPS (67 ms/frame); GPU batching keeps inference inside the live budget."
)
S7_PIPELINE = (
    "4× capture → YOLO ball + YOLO-Pose TRT FP16 → DLT/SVD + EMA/Kalman → "
    "ballistic solve + safety gates → USB serial → ESP32 FSM | "
    "Representative optimized log: total-loop P95 ≈64 ms"
)
S7_TAKEAWAY = (
    "Takeaway: the validated claim is sustained live-aim operation at 15 FPS; "
    "moving-target firing still depends on RPM-to-velocity calibration (Sec 6.4)."
)

def apply_s7(prs: _Prs) -> None:
    slide = prs.slides[6]
    set_paragraphs(find_shape_by_name(slide, "t101"), [S7_HEADLINE])
    set_paragraphs(find_shape_by_name(slide, "t103"), [S7_PIPELINE])
    set_paragraphs(find_shape_by_name(slide, "t106"), [S7_TAKEAWAY])

def check_s7(prs: _Prs) -> None:
    slide = prs.slides[6]
    assert_slide_text_contains(slide, "15 FPS (67 ms/frame)")
    assert_slide_text_contains(slide, "total-loop P95 ≈64 ms")
    assert_slide_text_not_contains(slide, "End-to-end runtime: ≈15 ms")
    assert_slide_text_not_contains(slide, "52 ms headroom")

# ---------- Slide 8: hardware — cost honesty ----------

S8_HEADLINE = "Costs: cameras ≈120, perception ≈200, BLM ≈358 USD."
S8_TAKEAWAY = (
    "Takeaway: expensive MoCap is not required for pose-guided aiming; "
    "the camera subset is ≈USD 120, thesis perception hardware is ≈USD 200, "
    "and the full BLM BOM is ≈USD 358 (Sec 1.3, 3.2)."
)
S8_AREA_TEXT_NEW = (
    "Arena setup: 4 USB cameras at corners, BLM at centre, "
    "24 AprilTag fiducials on walls for extrinsic calibration (Sec 3.5). "
    "All hardware fits in a domestic garage — no lab infrastructure required."
)

def apply_s8(prs: _Prs) -> None:
    slide = prs.slides[7]
    headline = find_shape_by_name(slide, "t101")
    headline.width = 6200000
    set_paragraphs(headline, [S8_HEADLINE])
    set_paragraphs(find_shape_by_name(slide, "t127"), [S8_TAKEAWAY])
    set_paragraphs(find_shape_by_name(slide, "t104"), [S8_AREA_TEXT_NEW])
    # The on-slide "cost table" is a grid of free-floating text shapes, not a
    # real GraphicFrame table. Update the last row (t124 "Perception PC" /
    # t125 "-") into a Total row so the deck has an explicit total line.
    set_paragraphs(find_shape_by_name(slide, "t124"), ["Total"])
    set_paragraphs(find_shape_by_name(slide, "t125"), ["≈USD 358"])

def check_s8(prs: _Prs) -> None:
    slide = prs.slides[7]
    assert_slide_text_contains(slide, "≈USD 358")
    assert_slide_text_contains(slide, "≈USD 120")
    assert_slide_text_contains(slide, "≈USD 200")
    assert_slide_text_contains(slide, "AprilTag fiducials")
    total_label = find_shape_by_name(slide, "t124")
    if shape_text(total_label).strip() != "Total":
        raise AssertionError(f"S8 grid 'Total' row mislabelled: {shape_text(total_label)!r}")

# ---------- Slide 9: software — de-clutter ----------

S9_FIRMWARE_EXCERPT = (
    "// control_12_full.ino\n"
    "void handle_set(){\n"
    "  // set <v> <h> <wl> <wr>\n"
    "  float v  = Serial.parseFloat();\n"
    "  float h  = Serial.parseFloat();\n"
    "  int   wl = Serial.parseInt();\n"
    "  int   wr = Serial.parseInt();\n"
    "  if(!armed) return;\n"
    "  clampAngle(v, h);\n"
    "  gateRPM(wl, wr);\n"
    "  target_v = v; target_h = h;\n"
    "}"
)

def apply_s9(prs: _Prs) -> None:
    slide = prs.slides[8]
    set_paragraphs(find_shape_by_name(slide, "t135"), [""])
    set_paragraphs(find_shape_by_name(slide, "t136"), [""])
    set_paragraphs(find_shape_by_name(slide, "t121"),
                   ["iterative reprojection-error rejection (Sec 3.7.2)"])
    set_paragraphs(find_shape_by_name(slide, "t125"),
                   ["adaptive EMA + CV Kalman (Sec 5.7)"])

def check_s9(prs: _Prs) -> None:
    slide = prs.slides[8]
    excerpt = find_shape_by_name(slide, "t136")
    if "handle_set" in shape_text(excerpt):
        raise AssertionError("S9 still carries the firmware code excerpt")
    assert_slide_text_contains(slide, "iterative reprojection-error rejection (Sec 3.7.2)")
    assert_slide_text_contains(slide, "adaptive EMA + CV Kalman (Sec 5.7)")

# ---------- New hidden appendix slide A5: firmware excerpt ----------

def apply_a5(prs: _Prs) -> None:
    """Create a hidden appendix slide with the firmware excerpt by cloning A3.
    Idempotent: only appends if no slide already has the A5 title."""
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and "A5 · Firmware command excerpt" in sh.text_frame.text:
                return
    from copy import deepcopy
    src = prs.slides[17]
    new_slide = prs.slides.add_slide(src.slide_layout)
    for ph in list(new_slide.shapes):
        ph._element.getparent().remove(ph._element)
    for sh in src.shapes:
        new_el = deepcopy(sh._element)
        new_slide.shapes._spTree.append(new_el)
    title = find_shape_by_name(new_slide, "Title 1")
    set_paragraphs(title, ["A5 · Firmware command excerpt"])
    set_paragraphs(find_shape_by_name(new_slide, "t101"), [
        "Firmware-level gates prevent stale or unsafe raw commands from bypassing the Python supervisor.",
    ])
    body = find_shape_by_name(new_slide, "t103")
    body.left = 685800
    body.top = 2133600
    body.width = 10820400
    body.height = 3657600
    set_paragraphs(body, S9_FIRMWARE_EXCERPT.split("\n"))
    set_paragraphs(find_shape_by_name(new_slide, "t105"), [
        "Use only if asked how the ESP32 enforces arm state, angle clamps, and RPM gates.",
    ])
    for sh in new_slide.shapes:
        if sh.has_text_frame and "[PLACEHOLDER_TABLE_1" in sh.text_frame.text:
            set_paragraphs(sh, [""])
    new_slide._element.set("show", "0")

def check_a5(prs: _Prs) -> None:
    found = None
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.has_text_frame and "A5 · Firmware command excerpt" in sh.text_frame.text:
                found = (i, s)
                break
        if found:
            break
    if not found:
        raise AssertionError("A5 firmware appendix slide missing")
    i, s = found
    if s._element.get("show") == "1":
        raise AssertionError(f"A5 (slide {i+1}) should be hidden")
    assert_slide_text_contains(s, "handle_set")
    assert_slide_text_contains(s, "control_12_full.ino")

# ---------- Slide 10: calibration — add numbers ----------

S10_BULLETS = [
    "Intrinsic: ChArUco board, reproj 2-8 px (Sec 5.1)",
    "Extrinsic: 24 AprilTags + RANSAC + σ=2.0 sigma-clipping (Sec 3.5.2)",
    "Extrinsic RMSE: 3-7 px after 8-15 % outlier rejection (Sec 5.2)",
    "Overlay validation: reprojected corners within 5-10 px on all 4 cams",
]

def apply_s10(prs: _Prs) -> None:
    slide = prs.slides[9]
    set_paragraphs(find_shape_by_name(slide, "t101"), S10_BULLETS)

def check_s10(prs: _Prs) -> None:
    slide = prs.slides[9]
    assert_slide_text_contains(slide, "2-8 px")
    assert_slide_text_contains(slide, "3-7 px")
    assert_slide_text_contains(slide, "5-10 px")

# ---------- Slide 11: methodology tightening ----------

S11_DYN_OUTCOME = "slow <800 mm jumps; fast blur stress; no-ball ≈0 FP"
S11_BIAS_OUTCOME = "150.77 → 95.17 mm (Fig 5.1)"
S11_NEW_TAKEAWAY = (
    "Takeaway: localisation, dynamic stability, and safety logging satisfy "
    "static + live-aim acceptance. Bias correction was fitted in-sample on "
    "the same 36-pt set (Sec 4.4.2, 6.3)."
)

def apply_s11(prs: _Prs) -> None:
    slide = prs.slides[10]
    set_paragraphs(find_shape_by_name(slide, "t133"), [S11_DYN_OUTCOME])
    set_paragraphs(find_shape_by_name(slide, "t141"), [S11_BIAS_OUTCOME])
    set_paragraphs(find_shape_by_name(slide, "t159"), [S11_NEW_TAKEAWAY])

def check_s11(prs: _Prs) -> None:
    slide = prs.slides[10]
    assert_slide_text_contains(slide, "150.77")
    assert_slide_text_contains(slide, "in-sample")
    assert_slide_text_not_contains(slide, "3D trajectory verified")

# ---------- Slide 12: key results — P95 + raw/corrected ----------

S12_SUBTITLE = (
    "All thresholds met within validated static/live-aim scope. "
    "Raw ball mean 150.77 mm → corrected 95.17 mm via in-sample bias model (Fig 5.1)."
)
S12_CAPTIONS = {
    "t105": "P95 166.51",
    "t107": "P95 198.73",
    "t109": "P95 171 / 172 / 200",
    "t111": "latch response",
    "t113": "6.2 ms/batch",
}

def apply_s12(prs: _Prs) -> None:
    slide = prs.slides[11]
    set_paragraphs(find_shape_by_name(slide, "t101"), [S12_SUBTITLE])
    for name, txt in S12_CAPTIONS.items():
        set_paragraphs(find_shape_by_name(slide, name), [txt])

def check_s12(prs: _Prs) -> None:
    slide = prs.slides[11]
    for needle in ["166.51", "198.73", "171 / 172 / 200", "150.77", "95.17"]:
        assert_slide_text_contains(slide, needle)
    assert_slide_text_contains(slide, "6.2 ms/batch")
    assert_slide_text_not_contains(slide, "52 ms headroom")

# ---------- Slide 13: limitations — expand to 6 bullets ----------

S13_LIMIT_LINES = [
    "Validated scope: static/live-aim single-shot tests; moving-subject firing remains future work.",
]
S13_NEW_TAKEAWAY = (
    "Takeaway: the limitations are explicit: moving-subject firing, in-sample "
    "bias fit, occlusion, and RPM-to-velocity calibration."
)

def apply_s13(prs: _Prs) -> None:
    slide = prs.slides[12]
    set_paragraphs(find_shape_by_name(slide, "t101"), S13_LIMIT_LINES)
    set_paragraphs(find_shape_by_name(slide, "t106"), [S13_NEW_TAKEAWAY])

def check_s13(prs: _Prs) -> None:
    slide = prs.slides[12]
    assert_slide_text_contains(slide, "moving-subject firing")
    assert_slide_text_contains(slide, "RPM-to-velocity calibration")

# ---------- Slide 14: conclusions — disadvantages + cost honesty ----------

S14_CARD_TEXT = {
    "TextBox 3":  "Live 3D joints\n→ BLM aim",
    "TextBox 4":  "Cameras ≈USD 120\nperception ≈USD 200\nBLM ≈USD 358",
    "TextBox 6":  "",
    "TextBox 7":  "Voice + keyboard\nbehind gates",
}
S14_NEW_TAKEAWAY = (
    "Takeaway: strong partial validation of a low-cost, pose-guided BLM. "
    "Disadvantages: in-sample bias fit · 3-camera occlusion floor · "
    "RPM → velocity not yet calibrated."
)

def apply_s14(prs: _Prs) -> None:
    slide = prs.slides[13]
    for name, txt in S14_CARD_TEXT.items():
        sh = find_shape_by_name(slide, name)
        set_paragraphs(sh, txt.split("\n"))
        set_text_size(sh, 13)
    set_paragraphs(find_shape_by_name(slide, "t121"), [S14_NEW_TAKEAWAY])

def check_s14(prs: _Prs) -> None:
    slide = prs.slides[13]
    assert_slide_text_contains(slide, "USD 358")
    assert_slide_text_contains(slide, "Disadvantages:")
    assert_slide_text_not_contains(slide, "closed-loop")

# ---------- Slide 15: future work + ethics + standards how ----------

S15_TAKEAWAY = (
    "Takeaway: moving-target work needs RPM calibration; Dual-use risk is gated "
    "by operator presence, NC E-STOP, and exclusion zone."
)
S15_STD_HOW = {
    "t107": "Machinery safety — L1–L10 hazard map (Sec 3.14)",
    "t110": "Safety-related control — NC E-STOP = ISO 13849-1 Cat-1 stop (L8)",
    "t113": "Wiring, fusing, E-STOP — 24V/50A fuse, single star-point ground (Sec 3.2)",
    "t116": "Operator-only zone during controlled fire (Sec 3.14.2)",
}

def apply_s15(prs: _Prs) -> None:
    slide = prs.slides[14]
    set_paragraphs(find_shape_by_name(slide, "t117"), [""])
    set_paragraphs(find_shape_by_name(slide, "t119"), [S15_TAKEAWAY])
    for name, txt in S15_STD_HOW.items():
        set_paragraphs(find_shape_by_name(slide, name), [txt])

def check_s15(prs: _Prs) -> None:
    slide = prs.slides[14]
    assert_slide_text_contains(slide, "Dual-use")
    assert_slide_text_contains(slide, "ISO 13849-1 Cat-1")
    assert_slide_text_contains(slide, "NC E-STOP")
    assert_slide_text_contains(slide, "24V/50A fuse")

# ---------- Slide 16: A1 live demo — embed video ----------

import os
import shutil as _shutil
import subprocess
from pathlib import Path as _Path

REPO_ROOT = _Path(__file__).resolve().parents[2]
DEMO_VIDEO = REPO_ROOT / "thesis_defense_presentation" / "IMG_1589 (online-video-cutter.com).mp4"
DEMO_POSTER = REPO_ROOT / "thesis_defense_presentation" / "_demo_poster.jpg"

def _ensure_poster() -> "_Path | None":
    if DEMO_POSTER.exists():
        return DEMO_POSTER
    if not _shutil.which("ffmpeg"):
        return None
    try:
        subprocess.run(
            ["ffmpeg", "-y", "-loglevel", "error", "-i", str(DEMO_VIDEO),
             "-vframes", "1", "-q:v", "3", str(DEMO_POSTER)],
            check=True, timeout=60,
        )
        return DEMO_POSTER
    except Exception:
        return None

def apply_s16(prs: _Prs) -> None:
    slide = prs.slides[15]
    placeholder = find_shape_by_name(slide, "t103")
    set_paragraphs(placeholder, [
        f"Embedded video: {DEMO_VIDEO.name}",
        "Duration: 54.5 s; play only the clearest 20-30 s segment if asked.",
        "Shows the live garage-arena system rather than a simulated pipeline.",
    ])
    from pptx.util import Inches
    placeholder.left = Inches(0.85)
    placeholder.top = Inches(2.05)
    placeholder.width = Inches(3.75)
    placeholder.height = Inches(3.15)
    set_paragraphs(find_shape_by_name(slide, "t101"), [
        "Appendix demo: real hardware video, cued only if the committee asks.",
    ])
    has_movie = any(
        getattr(sh, "shape_type", None) and "media" in str(sh.shape_type).lower()
        for sh in slide.shapes
    )
    if has_movie or not DEMO_VIDEO.exists():
        return
    poster = _ensure_poster()
    try:
        slide.shapes.add_movie(
            str(DEMO_VIDEO),
            left=Inches(5.30), top=Inches(1.55),
            width=Inches(2.70), height=Inches(4.80),
            poster_frame_image=str(poster) if poster else None,
            mime_type="video/mp4",
        )
    except Exception as e:
        set_paragraphs(placeholder,
                       [f"Demo clip: {DEMO_VIDEO.name} — drag in via LibreOffice ({e})"])

def check_s16(prs: _Prs) -> None:
    slide = prs.slides[15]
    assert_slide_text_not_contains(slide, "[PLACEHOLDER_VIDEO_1")
    assert_slide_text_contains(slide, "54.5 s")
    media_rels = [
        rel for rel in slide.part.rels.values()
        if "media" in rel.reltype or "video" in rel.reltype
    ]
    if not media_rels:
        raise AssertionError("S16 has no embedded video relationship")

# ---------- Slide 18: A3 latency table ----------

import json
import statistics as _stats

PERF_JSONL = REPO_ROOT / "Parallel_working" / "output" / "perf_blm_20260417_134342.jsonl"

LATENCY_FIELDS = ["capture_ms", "ball_ms", "pose_ms", "triang_ms",
                  "udp_ms", "viz3d_ms", "total_ms", "end_to_end_ms"]

def _build_latency_lines() -> list[str]:
    rows = [json.loads(l) for l in open(PERF_JSONL) if l.strip()]
    labels = {
        "capture_ms": "capture",
        "ball_ms": "ball detector",
        "pose_ms": "pose path",
        "triang_ms": "triangulation",
        "udp_ms": "UDP target",
        "viz3d_ms": "3D render",
        "total_ms": "total loop",
        "end_to_end_ms": "end-to-end diag.",
    }
    out = [f"Representative optimized log: {PERF_JSONL.name} ({len(rows)} frames)",
           "Stage              Mean ms   P95 ms"]
    for fld in LATENCY_FIELDS:
        vals = [r[fld] for r in rows if isinstance(r.get(fld), (int, float))]
        if not vals:
            continue
        mean = sum(vals) / len(vals)
        vals_sorted = sorted(vals)
        p95 = vals_sorted[int(0.95 * (len(vals_sorted) - 1))]
        out.append(f"{labels.get(fld, fld):<18}{mean:>7.1f}{p95:>9.1f}")
    return out

def apply_s18(prs: _Prs) -> None:
    slide = prs.slides[17]
    placeholder = find_shape_by_name(slide, "t103")
    if not PERF_JSONL.exists():
        set_paragraphs(placeholder, [f"Latency log not on disk: {PERF_JSONL.name}"])
        return
    set_paragraphs(placeholder, _build_latency_lines())
    set_paragraphs(find_shape_by_name(slide, "t105"), [
        "Takeaway: representative optimized live loop sustains 15 FPS "
        "(total-loop P95 ≈64 ms < 67 ms). Treat 6.2 ms as YOLO-Pose batch "
        "inference, not total end-to-end latency."
    ])

def check_s18(prs: _Prs) -> None:
    slide = prs.slides[17]
    assert_slide_text_not_contains(slide, "[PLACEHOLDER_TABLE_1")
    assert_slide_text_not_contains(slide, "Total perception: ~15 ms")
    assert_slide_text_contains(slide, "perf_blm_20260417_134342")
    assert_slide_text_contains(slide, "Mean ms")

# ---------- Registries ----------

APPLY_FUNCTIONS = [apply_s1, apply_s2, apply_s3, apply_s4, apply_s5, apply_s7,
                   apply_s8, apply_s9, apply_a5, apply_s10, apply_s11,
                   apply_s12, apply_s13, apply_s14, apply_s15,
                   apply_s16, apply_s18]
CHECK_FUNCTIONS = [check_s1, check_s2, check_s3, check_s4, check_s5, check_s7,
                   check_s8, check_s9, check_a5, check_s10, check_s11,
                   check_s12, check_s13, check_s14, check_s15,
                   check_s16, check_s18]
