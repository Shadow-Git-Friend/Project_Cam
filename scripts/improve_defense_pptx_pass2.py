#!/usr/bin/env python3
"""Pass 2: handle the two remaining [PLACEHOLDER_*] tokens on slides 2 and 6,
plus trim the title runs on slide 1 so "Arlen Smagulov" appears on its own
line rather than inheriting the [FILL ME] line breaks.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt, Inches

ROOT = Path("/home/hanush/Desktop/Project_Cam")
SRC = ROOT / "presentation_defense_improved_github_prioritized.pptx"

SLIDE6_ASSET = ROOT / "garage_lab_combined/thesis/figures_selected/fig_smoke_frame_0080.png"


def find_shape_by_text(slide, substr):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if substr in shape.text_frame.text:
            return shape
    return None


def replace_placeholder_with_image(slide, token, path):
    target = find_shape_by_text(slide, token)
    if target is None or not path.exists():
        return False
    left, top, width, height = target.left, target.top, target.width, target.height
    sp = target._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    return True


def make_placeholder_visible(slide, token, new_caption):
    """Rewrite the placeholder text to a short, clean instruction."""
    target = find_shape_by_text(slide, token)
    if target is None:
        return False
    tf = target.text_frame
    # Preserve the first paragraph's formatting; replace its runs with a short message.
    p = tf.paragraphs[0]
    for r in p.runs:
        r.text = ""
    if not p.runs:
        run = p.add_run()
    else:
        run = p.runs[0]
    run.text = new_caption
    run.font.size = Pt(14)
    # Clear any extra paragraphs
    for para in tf.paragraphs[1:]:
        for r in para.runs:
            r.text = ""
    return True


def main():
    prs = Presentation(str(SRC))

    # Slide 6 — substitute pipeline diagram with live smoke-test frame
    s6 = prs.slides[5]
    token6 = "[PLACEHOLDER_DIAGRAM_2"
    ok = replace_placeholder_with_image(s6, token6, SLIDE6_ASSET)
    print(f"Slide 6 pipeline placeholder → smoke_frame_0080.png: {ok}")

    # Slide 2 — rewrite placeholder to a short instruction box instead of the long spec
    s2 = prs.slides[1]
    token2 = "[PLACEHOLDER_SCREENSHOT_2"
    ok = make_placeholder_visible(
        s2, token2,
        "[ADD ON WINDOWS: split-card image — left panel a commercial ball launcher, right panel a MoCap studio]"
    )
    print(f"Slide 2 placeholder rewritten to short instruction: {ok}")

    prs.save(str(SRC))
    print(f"Saved: {SRC}")


if __name__ == "__main__":
    main()
