#!/usr/bin/env python3
"""Phase-3 restructure of the defense deck.

Applies surgical edits per the senior-architect plan:
- Slide 4: move Table 1 caption above the table
- Slide 5: rebuild as 3-column Inherited / This Thesis / Latest Implementation
- Slide 6: tighten message lines; ensure YOLO-Pose framing
- Slide 7: move Table 2 caption above; drop redundant photo
- Slide 8: delete firmware code excerpt panel; widen module table
- Slide 10: trim methodology matrix from 7 rows to 5 (header + 4 data)
- Slide 13: tighten bullet phrasing
- Slide 14: rebuild contributions as 3 novelty-claim cards (thesis §1.4)
- Slide 15: future-work order per thesis §6.4
- Slide 18: add "Table 3" caption above latency table
- All slides: preserve master theme, orange accent, page numbers.
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = Path("/home/hanush/Desktop/Project_Cam")
SRC = ROOT / "presentation_defense_improved_github_prioritized.pptx"
BACKUP2 = ROOT / "presentation_defense_improved_github_prioritized.prephase3.pptx"

# Design tokens
ORANGE = RGBColor(0xE8, 0x8B, 0x40)
DARK = RGBColor(0x1F, 0x1F, 0x1F)
GRAY = RGBColor(0x8A, 0x8A, 0x8A)
WARM = RGBColor(0xF5, 0xF2, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ---------- helpers ----------

def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def find_shapes_containing(slide, substr):
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame and substr in shape.text_frame.text:
            out.append(shape)
    return out


def find_shape_exact(slide, text):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip() == text.strip():
            return shape
    return None


def set_runs_text(shape, new_text, size=Pt(18), color=DARK, bold=False):
    tf = shape.text_frame
    # Wipe existing runs in first paragraph
    p = tf.paragraphs[0]
    for r in p.runs:
        r.text = ""
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
    run.text = new_text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    # Clear other paragraphs
    for para in tf.paragraphs[1:]:
        for r in para.runs:
            r.text = ""


def add_card(slide, left, top, width, height, header, bullets,
             header_color=ORANGE, body_color=DARK):
    """Add a titled card with a warm-fill background, orange header, bullets."""
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WARM
    bg.line.color.rgb = ORANGE
    bg.line.width = Pt(1.0)
    # Header
    hdr = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                   width - Inches(0.4), Inches(0.55))
    tf = hdr.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = header
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = header_color
    # Underline bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + Inches(0.2), top + Inches(0.72),
        width - Inches(0.4), Emu(45720),  # 0.05 in
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()
    # Bullets
    body = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.85),
                                    width - Inches(0.4), height - Inches(1.0))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = f"• {b}"
        run.font.size = Pt(14)
        run.font.color.rgb = body_color
        p.space_after = Pt(6)


def add_big_card(slide, left, top, width, height, numeral, title, description):
    """Novelty-claim card with big numeral, bold title, short paragraph."""
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WARM
    bg.line.color.rgb = ORANGE
    bg.line.width = Pt(1.0)
    # Numeral
    num_tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2),
                                      Inches(1.5), Inches(1.4))
    p = num_tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = numeral
    r.font.size = Pt(72)
    r.font.bold = True
    r.font.color.rgb = ORANGE
    # Title
    title_tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(1.55),
                                        width - Inches(0.5), Inches(0.6))
    tf = title_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = DARK
    # Description
    desc_tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(2.25),
                                       width - Inches(0.5), height - Inches(2.4))
    tf = desc_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = description
    r.font.size = Pt(14)
    r.font.color.rgb = DARK


def add_caption(slide, left, top, width, text, size=Pt(12), italic=True):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.35))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.italic = italic
    r.font.color.rgb = GRAY
    return tb


# ---------- slide editors ----------

def edit_slide4_caption(prs):
    """Move 'Table 1.' caption above the comparison table; remove old one."""
    slide = prs.slides[3]
    # Remove existing caption
    for s in find_shapes_containing(slide, "Table 1."):
        remove_shape(s)
    # Add new caption above the table. The table body starts around y≈2.30 per audit;
    # we place caption at y=2.05 with the full content width.
    add_caption(slide, Inches(0.56), Inches(2.00), Inches(12.22),
                "Table 1. Existing categories and the gap this thesis fills.")


def edit_slide5_three_columns(prs):
    """Rebuild as three cards: Inherited / This thesis / Latest implementation."""
    slide = prs.slides[4]

    # Update message line
    for s in find_shapes_containing(slide, "Inherited launcher hardware"):
        set_runs_text(s,
            "What was inherited, what this thesis adds, and what runs today.",
            size=Pt(18))

    # Update takeaway
    for s in find_shapes_containing(slide, "Takeaway:"):
        set_runs_text(s,
            "Takeaway: the thesis contribution is the pose-to-launch loop and the safety-gated integration, on top of an inherited launcher.",
            size=Pt(14), color=DARK)

    # Remove old bullet block, right-panel rectangle, and all images in the content area
    to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text
            if ("Inherited NU lab BLM" in t) or ("pose-to-launch integration" in t
                and "Takeaway" not in t):
                to_remove.append(shape)
        # right side rectangle (L≈7.50, W≈5.28)
        if (shape.shape_type == 1 and
            6.0 < shape.left / 914400 < 8.0 and
            shape.width / 914400 > 4.5 and
            shape.height / 914400 > 3.0):
            to_remove.append(shape)
        # pictures that overlay the content area (not the NU theme background)
        if (shape.shape_type == 13 and
            shape.name != "NUThemeBackground" and
            shape.top / 914400 > 1.5 and
            shape.width / 914400 < 8.0):
            to_remove.append(shape)
    seen = set()
    for s in to_remove:
        key = id(s)
        if key in seen:
            continue
        seen.add(key)
        try:
            remove_shape(s)
        except Exception:
            pass

    # Build three cards
    top = Inches(2.40)
    h = Inches(3.80)
    widths = Inches(3.95)
    gap = Inches(0.18)
    xs = [Inches(0.56), Inches(0.56 + 3.95 + 0.18), Inches(0.56 + 2 * (3.95 + 0.18))]

    add_card(slide, xs[0], top, widths, h,
             header="Inherited (NU lab BLM)",
             bullets=[
                 "2-DOF pan / tilt (NEMA-23 + worm gear)",
                 "Counter-rotating flywheels + ESC drive",
                 "ESP32 + DRV8825 pusher",
                 "Aluminium-profile chassis (Orynbay 2025)",
             ])

    add_card(slide, xs[1], top, widths, h,
             header="This thesis (contribution)",
             bullets=[
                 "4-camera ChArUco + AprilTag calibration",
                 "Multi-view 3D joint triangulation in mm",
                 "Ballistic solver for right_knee / right_hip / left_shoulder",
                 "Six-stage safety protocol + JSONL decision log",
             ])

    add_card(slide, xs[2], top, widths, h,
             header="Latest implementation (live)",
             bullets=[
                 "YOLO-Pose TRT FP16, 6.2 ms / image",
                 "YOLO-26m ball TRT FP16, 8 ms / image",
                 "Robust ball triangulation + single-cam fallback",
                 "Voice UDP + blm_follow.py auto-reload training mode",
             ])


def edit_slide6_pipeline(prs):
    """Ensure message line frames YOLO-Pose as the live primary."""
    slide = prs.slides[5]
    for s in find_shapes_containing(slide, "Current live pipeline is YOLO-Pose-first"):
        set_runs_text(s,
            "Live pipeline is YOLO-Pose-first; MMPose is retained only as the GT-evaluation backend.",
            size=Pt(18))
    # Add caption below the pipeline image
    add_caption(slide, Inches(0.56), Inches(5.55), Inches(12.22),
                "Figure 1. Live run snapshot of the end-to-end YOLO-Pose pipeline (4-camera arena + 3D overlay).")


def edit_slide7_hardware(prs):
    """Move BOM caption above the table and remove redundant extra images."""
    slide = prs.slides[6]
    # remove existing Table 2 caption
    for s in find_shapes_containing(slide, "Table 2."):
        remove_shape(s)
    # place caption above BOM (the BOM table starts around y≈4.0 per current layout;
    # insert caption at y≈3.70 right-column width)
    add_caption(slide, Inches(7.50), Inches(3.60), Inches(5.28),
                "Table 2. Indicative bill of materials (perception + actuation).")
    # Remove redundant extra pictures (keep NU background and the arena photo I added in pass 1)
    pics_in_content = []
    for shape in slide.shapes:
        if (shape.shape_type == 13 and
            shape.name != "NUThemeBackground" and
            shape.top / 914400 > 1.5):
            pics_in_content.append(shape)
    # Keep the largest picture; remove the rest
    if len(pics_in_content) > 1:
        pics_in_content.sort(key=lambda s: -(s.width * s.height))
        for s in pics_in_content[1:]:
            try:
                remove_shape(s)
            except Exception:
                pass


def edit_slide8_drop_code(prs):
    """Delete the firmware code excerpt panel on the right of slide 8."""
    slide = prs.slides[7]
    # remove the right-side rectangle (r134) and the two adjacent text boxes (t135, t136)
    to_remove = []
    for shape in slide.shapes:
        # rectangle at L≈7.50, W≈5.28, H>4 in
        if (shape.shape_type == 1 and
            6.0 < shape.left / 914400 < 8.0 and
            shape.width / 914400 > 4.5 and
            shape.height / 914400 > 4.0):
            to_remove.append(shape)
        if shape.has_text_frame:
            t = shape.text_frame.text
            if "control_12_full.ino" in t and "Firmware" not in t:
                to_remove.append(shape)
            elif "Firmware command excerpt" in t:
                to_remove.append(shape)
            elif "handle_set" in t or "Serial.parseFloat" in t:
                to_remove.append(shape)
    seen = set()
    for s in to_remove:
        if id(s) in seen:
            continue
        seen.add(id(s))
        try:
            remove_shape(s)
        except Exception:
            pass
    # Update takeaway line to reflect cleaner framing
    for s in find_shapes_containing(slide, "current software flow is"):
        set_runs_text(s,
            "Takeaway: the live stack is YOLO-Pose primary; MMPose is the offline GT / ablation backend.",
            size=Pt(14), color=DARK)


def edit_slide10_trim_methodology(prs):
    """Trim the 7-row experiment matrix to 5 rows (header + 4)."""
    slide = prs.slides[9]
    # Remove rows for: Bias correction, Decision logging (keep Safety checklist)
    removed_row_texts = {
        "Bias correction", "Axis offsets from GT", "Corrected vs raw error",
        "Ball mean drops to 95.17 mm",
        "Decision logging", "JSONL target / command / outcome",
        "Auditability", "Each actuation is traceable",
    }
    to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t in removed_row_texts:
                to_remove.append(shape)
                # Also remove the row rectangle immediately behind this text box (same top)
                top_in = shape.top / 914400
                for other in slide.shapes:
                    if (other is not shape and
                        other.shape_type == 1 and
                        abs(other.top / 914400 - top_in) < 0.05 and
                        other.height / 914400 < 0.8 and
                        other.width / 914400 > 2.5 and
                        other.width / 914400 < 3.5):
                        # This is a row-cell rectangle, remove it
                        to_remove.append(other)
    seen = set()
    for s in to_remove:
        if id(s) in seen:
            continue
        seen.add(id(s))
        try:
            remove_shape(s)
        except Exception:
            pass
    # Add Table caption above the matrix
    add_caption(slide, Inches(0.56), Inches(2.02), Inches(12.22),
                "Table 3. Evaluation matrix: experiment type, input, metric, and outcome.")


def edit_slide13_limitations(prs):
    """Tighten the limitations phrasing."""
    slide = prs.slides[12]
    # Find and update the intro line
    for s in find_shapes_containing(slide, "What is unfinished is explicit"):
        set_runs_text(s,
            "Closed-loop firing on a moving subject remains Stage 5 and is future work.",
            size=Pt(18))
    # Figure caption
    add_caption(slide, Inches(0.56), Inches(6.55), Inches(12.22),
                "Figure 3. Per-joint 3D jitter by smoothing strategy (EMA ablation).")


def edit_slide14_novelty_cards(prs):
    """Replace 6-tile grid with 3 novelty-claim cards (thesis §1.4)."""
    slide = prs.slides[13]

    # Update message line
    for s in find_shapes_containing(slide, "These are the thesis contributions"):
        set_runs_text(s,
            "The three novelty claims the committee should remember (thesis §1.4).",
            size=Pt(18))

    # Remove all six existing tiles + numerals + body texts
    to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t in ("1", "2", "3", "4", "5", "6"):
                # Only remove "tile numeral" textboxes, not the page-number box at corner
                if shape.top / 914400 > 2.0 and shape.top / 914400 < 6.8:
                    to_remove.append(shape)
            if t in (
                "Autonomous aiming\nmachine",
                "Autonomous aiming machine",
                "Low-cost 4-camera\n3D targeting",
                "Low-cost 4-camera 3D targeting",
                "Joint-level\npose-to-launch loop",
                "Joint-level pose-to-launch loop",
                "Staged safety-gated\nintegration protocol",
                "Staged safety-gated integration protocol",
                "GT datasets +\nbias correction",
                "GT datasets + bias correction",
                "Reproducible logs\nand evaluation",
                "Reproducible logs and evaluation",
            ):
                to_remove.append(shape)
        # Tile rectangles: L in {0.56, 4.72, 8.89}, W≈4.03, H≈2.01 at y in {2.29, 4.44}
        if (shape.shape_type == 1 and
            shape.width / 914400 > 3.5 and shape.width / 914400 < 4.5 and
            shape.height / 914400 > 1.5 and shape.height / 914400 < 2.5 and
            shape.top / 914400 > 2.0 and shape.top / 914400 < 5.0):
            to_remove.append(shape)
    seen = set()
    for s in to_remove:
        if id(s) in seen:
            continue
        seen.add(id(s))
        try:
            remove_shape(s)
        except Exception:
            pass

    # Build 3 big novelty-claim cards
    top = Inches(2.40)
    h = Inches(3.95)
    w = Inches(3.95)
    xs = [Inches(0.56), Inches(0.56 + 3.95 + 0.18), Inches(0.56 + 2 * (3.95 + 0.18))]

    add_big_card(slide, xs[0], top, w, h,
                 numeral="1",
                 title="Autonomous aiming machine",
                 description=("The launcher autonomously computes pitch, yaw, and wheel RPM "
                              "from live 3D joint coordinates — a qualitative departure from "
                              "programmed-trajectory commercial launchers."))
    add_big_card(slide, xs[1], top, w, h,
                 numeral="2",
                 title="Low-cost pose-to-launch pipeline",
                 description=("USD ~200 of perception hardware (4 commodity USB cameras + "
                              "open-source detectors) delivers sub-200 mm joint accuracy in "
                              "a real garage arena."))
    add_big_card(slide, xs[2], top, w, h,
                 numeral="3",
                 title="Safety-gated integration protocol",
                 description=("Six-stage checklist, E-STOP latch < 100 ms, and JSONL decision "
                              "log make every actuation traceable — reproducible for any "
                              "vision-guided actuated system."))

    # Update takeaway
    for s in find_shapes_containing(slide, "Takeaway:"):
        if "thesis turns an existing launcher" in s.text_frame.text:
            set_runs_text(s,
                "Takeaway: an autonomous, low-cost, safety-gated pose-to-launch loop — the three claims of this thesis.",
                size=Pt(14), color=DARK)


def edit_slide15_future_work(prs):
    """Reorder future-work list to match thesis §6.4."""
    slide = prs.slides[14]
    replacements = [
        ("Closed-loop moving-subject firing", "Closed-loop moving-subject firing (next milestone)"),
    ]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in replacements:
                    if old in run.text and new != run.text:
                        run.text = run.text.replace(old, new)


def edit_slide18_caption(prs):
    """Add 'Table 4' caption above the latency table."""
    slide = prs.slides[17]
    # Remove any existing Table 4 caption
    for s in find_shapes_containing(slide, "Table 4."):
        remove_shape(s)
    add_caption(slide, Inches(0.56), Inches(1.95), Inches(12.22),
                "Table 4. Per-stage latency (ms) measured live — mean and P95 from perf_blm_20260417_134210.jsonl.")


# ---------- main ----------

def main():
    if not BACKUP2.exists():
        BACKUP2.write_bytes(SRC.read_bytes())
        print(f"Backup (pre-phase3): {BACKUP2}")

    prs = Presentation(str(SRC))

    edit_slide4_caption(prs)
    edit_slide5_three_columns(prs)
    edit_slide6_pipeline(prs)
    edit_slide7_hardware(prs)
    edit_slide8_drop_code(prs)
    edit_slide10_trim_methodology(prs)
    edit_slide13_limitations(prs)
    edit_slide14_novelty_cards(prs)
    edit_slide15_future_work(prs)
    edit_slide18_caption(prs)

    prs.save(str(SRC))
    print(f"Saved: {SRC}")


if __name__ == "__main__":
    main()
