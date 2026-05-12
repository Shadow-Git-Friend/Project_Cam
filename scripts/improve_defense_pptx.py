#!/usr/bin/env python3
"""Improve the defense .pptx in-place: fix author/supervisor/date, replace
[PLACEHOLDER_*] tokens with real assets from the repo, align content with the
actual thesis by Arlen Smagulov. Preserves all style, color, and master-slide
design — only content and missing images are touched.
"""
from __future__ import annotations
import json
import statistics
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt, Inches

ROOT = Path("/home/hanush/Desktop/Project_Cam")
SRC = ROOT / "presentation_defense_improved_github_prioritized.pptx"
DST = ROOT / "presentation_defense_improved_github_prioritized.pptx"  # in-place
BACKUP = ROOT / "presentation_defense_improved_github_prioritized.backup.pptx"

# ---- Text replacements applied globally (run-level, style-preserving) ----
# Author/committee metadata
TEXT_REPLACEMENTS: list[tuple[str, str]] = [
    # Author name
    ("Hanush · MSc ECE · Nazarbayev University", "Arlen Smagulov · MSc ECE · Nazarbayev University"),
    ("Hanush, MSc ECE", "Arlen Smagulov, MSc ECE"),
    ("am Hanush, MSc ECE", "am Arlen Smagulov, MSc ECE"),  # speaker-notes opener
    # Title slide committee line
    ("Supervisor: [FILL ME]  ·  Committee: [FILL ME]  ·  Date: [FILL ME]",
     "Supervisor: Prof. Sultangali Arzykulov  ·  Co-supervisor: Prof. Mohammad Hashmi  ·  March 2026"),
    ("Supervisor: [FILL ME]", "Supervisor: Prof. Sultangali Arzykulov"),
    ("Committee: [FILL ME]", "Co-supervisor: Prof. Mohammad Hashmi"),
    ("Date: [FILL ME]", "March 2026"),
    # [VERIFY] tags whose citations are confirmed in the thesis bibliography
    ("Monocular pose (OpenPose [VERIFY], MediaPipe, YOLO-Pose)",
     "Monocular pose (OpenPose, MediaPipe, YOLO-Pose)"),
    ("Multi-view triangulation (Hartley & Zisserman, SVD-DLT [VERIFY])",
     "Multi-view triangulation (Hartley & Zisserman, SVD-DLT)"),
    # Title: use the actual thesis title. The pptx splits it across two lines.
    ("Pose-Guided Predictive Ballistics", "Pose-Guided Predictive Ballistics"),  # keep
    ("with Multi-Camera 3D Tracking", "for Body Part–Targeted Football Training"),
]

# Bare "Hanush" lone-run replacement (the title slide has "Hanush" in its own run)
LONE_HANUSH = ("Hanush", "Arlen Smagulov")

# Standalone replacement candidates (exact run equality)
EXACT_RUN_REPLACEMENTS: dict[str, str] = {
    "Hanush": "Arlen Smagulov",
}

# ---- Placeholder → asset mapping (slide_index_1based, placeholder_token, image_path) ----
# Each entry: (slide_idx_1based, placeholder_substring, asset_path_or_None)
PLACEHOLDER_ASSETS: list[tuple[int, str, Path | None]] = [
    # Slide 2: motivation split card — skip; no good split image exists
    (2, "[PLACEHOLDER_SCREENSHOT_2", None),
    # Slide 5: five-layer system stack diagram — use arena 360 fig as a hero
    (5, "[PLACEHOLDER_DIAGRAM_1", ROOT / "garage_lab_combined/thesis/figures_selected/fig_arena360_view_01.png"),
    # Slide 6: pipeline diagram — still missing
    (6, "[PLACEHOLDER_DIAGRAM_2", None),
    # Slide 7: arena hardware photo
    (7, "[PLACEHOLDER_SCREENSHOT_3", ROOT / "garage_lab_combined/thesis/figures_selected/fig_arena360_view_02.png"),
    # Slide 9: ChArUco + camera frusta — use calibration extrinsic overlay
    (9, "[PLACEHOLDER_DIAGRAM_3", ROOT / "garage_lab_combined/cal/extrinsics/arena_3d_view_1.png"),
    # Slide 11: key results — bias analysis + speed comparison
    (11, "[PLACEHOLDER_SCREENSHOT_P01", ROOT / "Parallel_working/output/ablation_results/viz_gt_bias_analysis.png"),
    (11, "[PLACEHOLDER_SCREENSHOT_P05", ROOT / "Parallel_working/output/ablation_results/viz_speed_comparison.png"),
    # Slide 12: analysis — backend compare + joint errors
    (12, "[PLACEHOLDER_SCREENSHOT_P04", ROOT / "Parallel_working/output/ablation_results/viz_backend_comparison.png"),
    (12, "[PLACEHOLDER_SCREENSHOT_P02", ROOT / "garage_lab_combined/thesis/figures_selected/fig_joint_touch_error_boxplot.png"),
    # Slide 13: limitations — EMA ablation
    (13, "[PLACEHOLDER_SCREENSHOT_P06", ROOT / "Parallel_working/output/ablation_results/viz_ema_ablation_jitter.png"),
    # Slide 16: demo — we embed a still frame since pptx video embed is heavy/fragile.
    # Use a smoke-test frame as the poster image with a caption pointing to the video path.
    (16, "[PLACEHOLDER_VIDEO_1", ROOT / "garage_lab_combined/thesis/figures_selected/fig_smoke_frame_0200.png"),
    # Slide 18: latency table — built procedurally below
    (18, "[PLACEHOLDER_TABLE_1", None),  # handled specially
]

PERF_JSONL = ROOT / "Parallel_working/output/perf_blm_20260417_134210.jsonl"

# ----------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------

def replace_text_in_run(run, old: str, new: str) -> bool:
    if old in run.text:
        run.text = run.text.replace(old, new)
        return True
    return False


def apply_global_text_replacements(prs: Presentation) -> int:
    n = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in TEXT_REPLACEMENTS:
                        if replace_text_in_run(run, old, new):
                            n += 1
                    # Exact-run replacements (whole run equals old)
                    if run.text.strip() in EXACT_RUN_REPLACEMENTS:
                        run.text = EXACT_RUN_REPLACEMENTS[run.text.strip()]
                        n += 1
        # also update notes
        if slide.has_notes_slide:
            for para in slide.notes_slide.notes_text_frame.paragraphs:
                for run in para.runs:
                    for old, new in TEXT_REPLACEMENTS:
                        replace_text_in_run(run, old, new)
    return n


def find_shape_by_text(slide, substr: str):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if substr in shape.text_frame.text:
            return shape
    return None


def replace_placeholder_with_image(slide, placeholder_substr: str, image_path: Path) -> bool:
    target = find_shape_by_text(slide, placeholder_substr)
    if target is None:
        return False
    if not image_path.exists():
        print(f"  ! image missing: {image_path}")
        return False
    left, top, width, height = target.left, target.top, target.width, target.height
    # Remove placeholder shape
    sp = target._element
    sp.getparent().remove(sp)
    # Insert picture with same bounding box; pptx preserves aspect by default if one dim omitted
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    return True


def compute_latency_stats(jsonl_path: Path) -> dict[str, tuple[float, float]]:
    """Return {stage: (mean_ms, p95_ms)} from perf JSONL."""
    stages = ["capture_ms", "ball_ms", "pose_ms", "triang_ms", "udp_ms",
              "viz3d_ms", "mosaic_ms", "end_to_end_ms", "total_ms"]
    data: dict[str, list[float]] = {s: [] for s in stages}
    with jsonl_path.open() as f:
        for line in f:
            if not line.strip():
                continue
            rec = json.loads(line)
            for s in stages:
                v = rec.get(s)
                if v is not None:
                    data[s].append(float(v))
    out: dict[str, tuple[float, float]] = {}
    for s, vals in data.items():
        if not vals:
            continue
        vals_sorted = sorted(vals)
        p95 = vals_sorted[int(0.95 * (len(vals_sorted) - 1))]
        out[s] = (statistics.mean(vals), p95)
    return out


LATENCY_STAGE_LABELS: list[tuple[str, str]] = [
    ("capture_ms", "Camera capture (per-frame)"),
    ("ball_ms", "YOLO ball detection"),
    ("pose_ms", "MMPose / YOLO-Pose"),
    ("triang_ms", "SVD-DLT triangulation"),
    ("udp_ms", "UDP target broadcast"),
    ("viz3d_ms", "3D render (cv2 backend)"),
    ("mosaic_ms", "2D mosaic render"),
    ("end_to_end_ms", "End-to-end perception"),
    ("total_ms", "Total loop (incl. render)"),
]


def insert_latency_table(slide, placeholder_substr: str, stats: dict[str, tuple[float, float]]) -> bool:
    target = find_shape_by_text(slide, placeholder_substr)
    if target is None:
        return False
    left, top, width, height = target.left, target.top, target.width, target.height
    sp = target._element
    sp.getparent().remove(sp)
    rows = 1 + sum(1 for k, _ in LATENCY_STAGE_LABELS if k in stats)
    cols = 3
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    # Header
    headers = ["Stage", "Mean (ms)", "P95 (ms)"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(14)
    r = 1
    for key, label in LATENCY_STAGE_LABELS:
        if key not in stats:
            continue
        mean_ms, p95_ms = stats[key]
        row_vals = [label, f"{mean_ms:.1f}", f"{p95_ms:.1f}"]
        for j, v in enumerate(row_vals):
            cell = table.cell(r, j)
            cell.text = v
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)
        r += 1
    return True


# ---- Targeted text edits on specific slides (content alignment) ----

SLIDE_EDITS: dict[int, list[tuple[str, str]]] = {
    # Slide 2 motivation — emphasize football training explicitly
    2: [
        ("Coaches need joint-specific delivery for training and rehab.",
         "Football coaches need joint-specific delivery for reception, reaction, and rehab drills."),
    ],
    # Slide 4 background — mention Yessimkhan precursor and football scope
    4: [
        ("Research gap: low-cost joint targeting with physical actuation.",
         "Gap: low-cost joint-level targeting with physical actuation. Prior NU work (Orynbay 2025) delivered an omnidirectional launcher but no joint-level vision guidance."),
    ],
    # Slide 5 — clarify scope matches Arlen's thesis
    5: [
        ("Inherited platform: 2-DOF BLM, flywheels, ESP32, and stepper-driven pan/tilt.",
         "Inherited NU lab BLM: 2-DOF pan/tilt, counter-rotating flywheels, ESP32 low-level control."),
        ("This thesis adds 4-camera perception, 3D triangulation, ballistic targeting, and safety gating.",
         "This thesis adds 4-camera perception, mm-accurate 3D triangulation, ballistic targeting of right_knee/right_hip/left_shoulder, and six-stage safety gating."),
    ],
    # Slide 10 methodology — replace ambiguous dynamic clip with thesis wording
    10: [
        ("ball_slow / ball_fast / no_ball", "ball_slow / ball_fast / ball_fast_ema0.1 / no_ball"),
        ("Aim-only + static shot PASS", "S0–S4 PASSED (aim-only + controlled static single-shot)"),
    ],
    # Slide 13 limitations — add calibration single-session in-sample caveat from thesis §4.4.2
    13: [
        ("Pan / tilt homing and long-session drift robustness need more validation.",
         "Bias correction fit in-sample on the 36-point ball grid; no held-out calibration set."),
    ],
    # Slide 14 contributions — align wording to Arlen's three stated novelty claims
    14: [
        ("Autonomous aiming", "Autonomous aiming"),  # keep
        ("YOLO-Pose to", "Joint-level"),
        ("launch control loop", "pose-to-launch loop"),
        ("Six-stage safety", "Staged safety-gated"),
        ("validation", "integration protocol"),
    ],
    # Slide 15 future work — reorder to match Arlen's §6.4 sections
    15: [
        ("Closed-loop moving-subject firing", "Closed-loop autonomous firing on moving subjects"),
        ("Empirical ballistic calibration map", "Empirical RPM → m/s ballistic calibration map"),
        ("SLAM-based camera relocalisation", "SLAM-based camera re-localisation and self-recalibration"),
        ("Virtual 3D Goal for camera-only impact detection",
         "Virtual 3D Goal: camera-only impact detection, no sensors at the goal"),
    ],
}


def apply_slide_edits(prs: Presentation) -> int:
    n = 0
    for idx, edits in SLIDE_EDITS.items():
        slide = prs.slides[idx - 1]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in edits:
                        if old == new:
                            continue
                        if replace_text_in_run(run, old, new):
                            n += 1
    return n


# ----------------------------------------------------------------------
# Main
# ----------------------------------------------------------------------

def main():
    # Back up once
    if not BACKUP.exists():
        BACKUP.write_bytes(SRC.read_bytes())
        print(f"Backup: {BACKUP}")

    prs = Presentation(str(SRC))
    print(f"Loaded: {len(prs.slides)} slides")

    # 1) Global text replacements (author, supervisor, committee, date, [VERIFY])
    n = apply_global_text_replacements(prs)
    print(f"Text replacements applied: {n}")

    # 2) Slide-specific content corrections
    n = apply_slide_edits(prs)
    print(f"Slide-specific edits: {n}")

    # 3) Placeholder → image replacements
    for idx_1b, token, path in PLACEHOLDER_ASSETS:
        slide = prs.slides[idx_1b - 1]
        if token == "[PLACEHOLDER_TABLE_1":
            stats = compute_latency_stats(PERF_JSONL)
            ok = insert_latency_table(slide, token, stats)
            print(f"  slide {idx_1b} table {'OK' if ok else 'NOT FOUND'}")
            continue
        if path is None:
            print(f"  slide {idx_1b} {token} left as placeholder (no asset)")
            continue
        ok = replace_placeholder_with_image(slide, token, path)
        print(f"  slide {idx_1b} {token[:40]} {'→ ' + path.name if ok else 'NOT FOUND'}")

    # 4) Add a visible caption on slide 16 pointing to the real video path
    slide16 = prs.slides[15]
    from pptx.util import Inches as In
    caption_left = In(0.5)
    caption_top = In(6.5)
    caption_w = In(12.3)
    caption_h = In(0.5)
    tb = slide16.shapes.add_textbox(caption_left, caption_top, caption_w, caption_h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = ("Poster frame shown above. Live demo video on disk: "
              "Parallel_working/output/recordings/arena3d_20260417_123348.mp4 "
              "(trim to 20–30 s, play on demand).")
    r.font.size = Pt(12)

    prs.save(str(DST))
    print(f"Saved: {DST}")


if __name__ == "__main__":
    main()
