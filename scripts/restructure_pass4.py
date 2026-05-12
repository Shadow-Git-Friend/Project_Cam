#!/usr/bin/env python3
"""Pass 4: fix the Firmware row on slide 8 and add figure captions on slides 11 and 12."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path("/home/hanush/Desktop/Project_Cam")
SRC = ROOT / "presentation_defense_improved_github_prioritized.pptx"

GRAY = RGBColor(0x8A, 0x8A, 0x8A)
DARK = RGBColor(0x1F, 0x1F, 0x1F)


def add_caption(slide, left, top, width, text, size=Pt(12)):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.32))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.italic = True
    r.font.color.rgb = GRAY


def main():
    prs = Presentation(str(SRC))

    # --- Slide 8: re-insert the Firmware row module/tool cell ---
    s8 = prs.slides[7]
    tb = s8.shapes.add_textbox(Inches(2.61), Inches(6.05), Inches(4.36), Inches(0.55))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "control_12_full.ino (cooperative FSM, 921 600 baud)"
    r.font.size = Pt(14)
    r.font.color.rgb = DARK

    # --- Slide 11: Figure 4 + Figure 5 captions ---
    # Inspect actual image geometry first
    s11 = prs.slides[10]
    pics = [sh for sh in s11.shapes if sh.shape_type == 13 and sh.name != "NUThemeBackground"
            and sh.top / 914400 > 2.0]
    pics.sort(key=lambda p: p.left)
    if len(pics) >= 2:
        left_pic = pics[0]
        right_pic = pics[1]
        # Place caption directly under each picture
        add_caption(s11, left_pic.left, left_pic.top + left_pic.height,
                    left_pic.width,
                    "Figure 4. Ball static localisation: raw vs bias-corrected 3D error (95.17 mm mean corrected).")
        add_caption(s11, right_pic.left, right_pic.top + right_pic.height,
                    right_pic.width,
                    "Figure 5. Per-backend latency: YOLO-Pose (TRT FP16) is 6.2× faster than MMPose.")

    # --- Slide 12: Figure 6 + Figure 7 captions ---
    s12 = prs.slides[11]
    pics = [sh for sh in s12.shapes if sh.shape_type == 13 and sh.name != "NUThemeBackground"
            and sh.top / 914400 > 2.0]
    pics.sort(key=lambda p: p.left)
    if len(pics) >= 2:
        left_pic = pics[0]
        right_pic = pics[1]
        add_caption(s12, left_pic.left, left_pic.top + left_pic.height,
                    left_pic.width,
                    "Figure 6. Backend comparison: YOLO-Pose matches MMPose 3D jitter within 5 mm.")
        add_caption(s12, right_pic.left, right_pic.top + right_pic.height,
                    right_pic.width,
                    "Figure 7. Per-joint 3D error increases with height (knee → hip → shoulder) due to visibility geometry.")

    prs.save(str(SRC))
    print(f"Saved: {SRC}")


if __name__ == "__main__":
    main()
