#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parent
MEDIA = ROOT / "media"
OUT_PPTX = ROOT / "kairat_academy_pilot_pitch.pptx"
OUT_NOTES = ROOT / "speaker_notes_ru.md"

W = Cm(33.867)
H = Cm(19.05)

BLACK = RGBColor(18, 18, 18)
DARK = RGBColor(31, 31, 31)
GRAY = RGBColor(96, 96, 96)
LIGHT_GRAY = RGBColor(244, 244, 242)
MID_GRAY = RGBColor(220, 220, 215)
YELLOW = RGBColor(255, 211, 0)
GOLD = RGBColor(226, 180, 0)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(72, 170, 86)
RED = RGBColor(205, 69, 57)


def slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(s, x, y, w, h, fill, line=None, radius=False):
    shape = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        x,
        y,
        w,
        h,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line if line is not None else fill
    shape.line.width = Pt(0.5)
    return shape


def textbox(s, x, y, w, h, text, size=24, color=BLACK, bold=False, align=PP_ALIGN.LEFT, font="Arial"):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Cm(0.05)
    tf.margin_right = Cm(0.05)
    tf.margin_top = Cm(0.03)
    tf.margin_bottom = Cm(0.03)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def title(s, text, subtitle=None):
    textbox(s, Cm(1.05), Cm(0.72), Cm(25.5), Cm(1.15), text, 31, BLACK, True)
    rect(s, Cm(1.05), Cm(1.95), Cm(5.9), Cm(0.12), YELLOW)
    if subtitle:
        textbox(s, Cm(1.05), Cm(2.25), Cm(25.7), Cm(0.7), subtitle, 12.5, GRAY)


def footer(s, n):
    textbox(s, Cm(1.05), Cm(18.25), Cm(10), Cm(0.35), "Project_Cam × Kairat Academy", 8.5, GRAY)
    textbox(s, Cm(31.65), Cm(18.25), Cm(0.8), Cm(0.35), str(n), 8.5, GRAY, align=PP_ALIGN.RIGHT)


def bullet(s, x, y, w, text, accent=YELLOW, size=16.5):
    rect(s, x, y + Cm(0.13), Cm(0.13), Cm(0.55), accent)
    textbox(s, x + Cm(0.35), y, w - Cm(0.35), Cm(0.82), text, size, BLACK)


def metric(s, x, y, label, value, note, color=BLACK, value_size=None):
    if value_size is None:
        plain_len = len(value.replace("\n", ""))
        value_size = 24 if plain_len <= 8 else 20 if plain_len <= 12 else 17
    rect(s, x, y, Cm(5.2), Cm(2.25), WHITE, MID_GRAY, radius=True)
    textbox(s, x + Cm(0.28), y + Cm(0.25), Cm(4.6), Cm(0.42), label, 9.5, GRAY, True)
    textbox(s, x + Cm(0.28), y + Cm(0.72), Cm(4.6), Cm(0.72), value, value_size, color, True)
    textbox(s, x + Cm(0.28), y + Cm(1.48), Cm(4.55), Cm(0.55), note, 8.8, GRAY)


def add_video_panel(s, video_path, poster_path, x, y, w, h):
    s.shapes.add_movie(
        str(video_path),
        x,
        y,
        w,
        h,
        poster_frame_image=str(poster_path),
        mime_type="video/mp4",
    )
    rect(s, x + w - Cm(2.0), y + h - Cm(1.35), Cm(1.35), Cm(0.95), BLACK, BLACK, radius=True)
    tri = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE, x + w - Cm(1.58), y + h - Cm(1.14), Cm(0.48), Cm(0.48))
    tri.fill.solid()
    tri.fill.fore_color.rgb = YELLOW
    tri.line.color.rgb = YELLOW


def slide_1(prs):
    s = slide(prs)
    rect(s, 0, 0, W, H, LIGHT_GRAY)
    img = MEDIA / "live_blm_aim_demo_poster.jpg"
    s.shapes.add_picture(str(img), Cm(18.15), Cm(0), Cm(15.72), H)
    rect(s, Cm(17.5), 0, Cm(1.2), H, YELLOW)
    textbox(s, Cm(1.15), Cm(0.92), Cm(8.4), Cm(0.42), "PROJECT_CAM / PROXIBALL 3D", 10.5, GRAY, True)
    textbox(s, Cm(1.1), Cm(2.15), Cm(15.5), Cm(2.6), "Проект для\nАкадемии Кайрат", 34, BLACK, True)
    textbox(
        s,
        Cm(1.15),
        Cm(5.45),
        Cm(14.2),
        Cm(2.05),
        "Мы превращаем обычную тренировку в измеряемую, персональную и интерактивную.",
        19.5,
        DARK,
        True,
    )
    textbox(s, Cm(1.18), Cm(8.02), Cm(13.2), Cm(0.8), "3D-поза игрока + умная подача мяча + безопасная роботизированная пушка", 14.5, GRAY)
    metric(s, Cm(1.15), Cm(10.45), "Цель", "3 мес.", "пилот в академии", YELLOW)
    metric(s, Cm(6.75), Cm(10.45), "Фокус", "U-команды", "дети и юниоры", BLACK)
    metric(s, Cm(12.35), Cm(10.45), "Формат", "1 арена", "установка + отчёты", BLACK)
    textbox(s, Cm(1.15), Cm(16.15), Cm(12.6), Cm(0.75), "Для тренера: понятная обратная связь по реакции, движению и точности без лабораторной инфраструктуры.", 14, BLACK)
    footer(s, 1)


def slide_2(prs):
    s = slide(prs)
    rect(s, 0, 0, W, H, LIGHT_GRAY)
    title(s, "Что уже работает", "Показываем не идею, а живой прототип: человек, 3D-скелет, прицеливание и физическая пушка.")
    add_video_panel(
        s,
        MEDIA / "live_blm_aim_demo.mp4",
        MEDIA / "live_blm_aim_demo_poster.jpg",
        Cm(1.05),
        Cm(3.35),
        Cm(17.0),
        Cm(12.2),
    )
    bullet(s, Cm(19.1), Cm(3.55), Cm(12.4), "4 камеры строят 3D-положение игрока в реальном времени.", YELLOW, 15.2)
    bullet(s, Cm(19.1), Cm(4.95), Cm(12.4), "Система выбирает сустав: колено, плечо, корпус или голову.", YELLOW, 15.2)
    bullet(s, Cm(19.1), Cm(6.35), Cm(12.4), "BLM получает безопасную команду: угол, направление и скорость.", YELLOW, 15.2)
    bullet(s, Cm(19.1), Cm(7.75), Cm(12.4), "Цикл reload → aim → shoot уже проверен на стенде.", YELLOW, 15.2)
    metric(s, Cm(19.15), Cm(9.2), "Живой режим", "15 FPS", "реальная частота прототипа", GREEN)
    metric(s, Cm(24.75), Cm(9.2), "Задержка", "~64 ms", "средний live-loop", GREEN)
    metric(s, Cm(19.15), Cm(12.05), "Безопасность", "S0–S4", "поэтапная проверка BLM", GREEN)
    metric(s, Cm(24.75), Cm(12.05), "Управление", "3D→aim", "3D цель → команда", GREEN)
    textbox(s, Cm(19.15), Cm(15.35), Cm(12.0), Cm(1.25), "Главная мысль: прототип уже видит игрока, понимает цель и переводит её в действие машины.", 15, BLACK, True)
    footer(s, 2)


def slide_3(prs):
    s = slide(prs)
    rect(s, 0, 0, W, H, LIGHT_GRAY)
    title(s, "Как это помогает тренеру", "Система не заменяет тренера. Она даёт измерения, которые тренер обычно видит только субъективно.")
    add_video_panel(
        s,
        MEDIA / "arena_3d_skeleton.mp4",
        MEDIA / "arena_3d_skeleton_poster.jpg",
        Cm(1.05),
        Cm(3.35),
        Cm(13.35),
        Cm(12.55),
    )
    items = [
        ("Тест реакции вратаря", "Пушка подаёт в зону, система считает время реакции и движение центра тела."),
        ("Адресные упражнения", "Можно тренировать приём, уход, блок и координацию через конкретные зоны тела."),
        ("Оценка движения", "3D-скелет помогает увидеть асимметрию, колено, корпус и качество повторов."),
        ("Отчёт после сессии", "Тренер и физиотерапевт получают понятный файл: что улучшать и что повторить."),
    ]
    y = Cm(3.55)
    for head, body in items:
        rect(s, Cm(15.25), y, Cm(16.8), Cm(2.35), WHITE, MID_GRAY, radius=True)
        textbox(s, Cm(15.75), y + Cm(0.32), Cm(15.4), Cm(0.48), head, 15.5, BLACK, True)
        textbox(s, Cm(15.75), y + Cm(0.93), Cm(15.2), Cm(0.95), body, 12.2, GRAY)
        y += Cm(2.75)
    textbox(s, Cm(15.35), Cm(15.35), Cm(15.8), Cm(1.0), "Для академии это масштабируемый формат: один короткий тест можно повторять для десятков игроков и сравнивать прогресс.", 14.5, BLACK, True)
    footer(s, 3)


def slide_4(prs):
    s = slide(prs)
    rect(s, 0, 0, W, H, LIGHT_GRAY)
    title(s, "Почему это отличается от обычной пушки", "Ценность не в механике отдельно, а в связке: измерение игрока → адаптивная подача → отчёт для тренера.")
    rows = [
        ("Обычная пушка", "Повторяет один и тот же удар", "Дёшево, но не реагирует на игрока"),
        ("Motion-capture лаборатория", "Очень точные измерения движения", "Дорого, маркеры, не подходит для ежедневной тренировки"),
        ("Project_Cam", "3D-трекинг + физическая подача + безопасность", "Пилотируемая система для футбольной академии"),
    ]
    x0, y0 = Cm(1.05), Cm(3.5)
    col_w = [Cm(7.2), Cm(11.7), Cm(11.7)]
    headers = ["Подход", "Что умеет", "Ограничение / преимущество"]
    x = x0
    for i, htxt in enumerate(headers):
        rect(s, x, y0, col_w[i], Cm(1.0), BLACK if i == 0 else DARK)
        textbox(s, x + Cm(0.25), y0 + Cm(0.22), col_w[i] - Cm(0.5), Cm(0.5), htxt, 12.5, WHITE, True)
        x += col_w[i]
    y = y0 + Cm(1.0)
    for idx, row in enumerate(rows):
        row_fill = WHITE if idx < 2 else RGBColor(255, 248, 214)
        x = x0
        for i, cell in enumerate(row):
            rect(s, x, y, col_w[i], Cm(2.15), row_fill, MID_GRAY)
            textbox(s, x + Cm(0.25), y + Cm(0.35), col_w[i] - Cm(0.5), Cm(1.3), cell, 13.2 if i else 13.7, BLACK, i == 0 or idx == 2)
            x += col_w[i]
        y += Cm(2.15)
    textbox(s, Cm(1.05), Cm(11.75), Cm(30.0), Cm(0.7), "Доказательная база прототипа", 18, BLACK, True)
    metric(s, Cm(1.05), Cm(12.85), "Мяч", "95.17 mm", "средняя ошибка после коррекции", GREEN)
    metric(s, Cm(6.65), Cm(12.85), "Суставы", "143.38 mm", "средняя ошибка joint-touch", GREEN)
    metric(s, Cm(12.25), Cm(12.85), "Скорость", "15 FPS", "живой режим", GREEN)
    metric(s, Cm(17.85), Cm(12.85), "Безопасность", "10 слоёв", "ПО + прошивка + оператор", GREEN)
    metric(s, Cm(23.45), Cm(12.85), "Демо", "видео", "реальное железо", GREEN)
    footer(s, 4)


def slide_5(prs):
    s = slide(prs)
    rect(s, 0, 0, W, H, LIGHT_GRAY)
    title(s, "Пилот с Кайратом: что нужно", "Предложение: 3-месячный пилот, чтобы проверить систему на реальных тренировочных сценариях академии.")
    rect(s, Cm(1.05), Cm(3.35), Cm(13.7), Cm(11.3), BLACK, BLACK, radius=True)
    textbox(s, Cm(1.75), Cm(4.05), Cm(11.9), Cm(0.65), "Инвестиционный запрос", 14, YELLOW, True)
    textbox(s, Cm(1.75), Cm(5.05), Cm(11.5), Cm(1.95), "3-месячный\nпилот академии", 27, WHITE, True)
    textbox(s, Cm(1.75), Cm(7.45), Cm(11.6), Cm(1.25), "Первый конкретный пункт: апгрейд камер и синхронизации для стабильной работы на скорости.", 12.8, WHITE)
    textbox(s, Cm(1.75), Cm(9.25), Cm(11.6), Cm(0.62), "~1.84 млн KZT на оборудование", 18.5, YELLOW, True)
    textbox(s, Cm(1.75), Cm(10.12), Cm(11.3), Cm(1.05), "6 global-shutter камер, 60 FPS, hardware trigger, GigE, объективы, NVMe и синхронизация.", 11.5, WHITE)
    textbox(s, Cm(1.75), Cm(12.25), Cm(11.7), Cm(1.25), "Это не «починка прототипа», а переход от лабораторного proof-of-concept к надёжности академического масштаба.", 12.3, YELLOW, True)
    months = [
        ("Месяц 1", "установка, калибровка, safety check"),
        ("Месяц 2", "тренерские сессии: реакция, зоны, отчёты"),
        ("Месяц 3", "итоговый отчёт и решение о внедрении"),
    ]
    y = Cm(3.55)
    for head, body in months:
        rect(s, Cm(16.0), y, Cm(15.6), Cm(2.5), WHITE, MID_GRAY, radius=True)
        textbox(s, Cm(16.55), y + Cm(0.38), Cm(3.4), Cm(0.5), head, 15.5, BLACK, True)
        rect(s, Cm(20.15), y + Cm(0.35), Cm(0.12), Cm(1.65), YELLOW)
        textbox(s, Cm(20.55), y + Cm(0.35), Cm(9.9), Cm(1.1), body, 15, DARK)
        y += Cm(2.95)
    rect(s, Cm(16.0), Cm(12.85), Cm(15.6), Cm(2.65), RGBColor(255, 248, 214), GOLD, radius=True)
    textbox(s, Cm(16.55), Cm(13.25), Cm(14.2), Cm(0.55), "Что получает Кайрат", 16, BLACK, True)
    textbox(s, Cm(16.55), Cm(14.0), Cm(14.1), Cm(0.95), "пилотную спортивную технологию, данные по игрокам, видео-доказательство и понятный план масштабирования.", 13.5, DARK)
    textbox(s, Cm(16.0), Cm(16.6), Cm(14.4), Cm(0.7), "Источник масштаба: официальный сайт ФК «Кайрат» указывает 1000+ воспитанников и 50+ специалистов академии.", 10.5, GRAY)
    footer(s, 5)


def build_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)
    prs.save(OUT_PPTX)


def build_notes():
    notes = """# Speaker Notes — Kairat Academy Pilot Pitch

## Slide 1 — Проект для Академии Кайрат
Коротко открыть: это не академическая защита, а предложение пилота для футбольной академии. Главная фраза: мы превращаем обычную тренировку в измеряемую, персональную и интерактивную. Не объяснять алгоритмы; сказать простыми словами: камеры видят игрока в 3D, система выбирает цель, пушка подаёт мяч безопасно и управляемо.

## Slide 2 — Что уже работает
Показать видео. Сказать, что это реальное железо: 4 камеры, 3D-скелет, выбранный сустав, команда на BLM. Подчеркнуть, что прототип уже проходит путь от позы игрока до команды прицеливания. Цифры держать коротко: 15 FPS, около 64 ms live-loop, S0-S4 проверки безопасности.

## Slide 3 — Как это помогает тренеру
Показать 3D-скелет. Объяснить ценность для тренера: реакция вратаря, адресные упражнения, контроль движения, отчёт для тренера и физиотерапевта. Важно: система не заменяет тренера, она делает тренировку измеряемой и повторяемой.

## Slide 4 — Почему это отличается от обычной пушки
Сравнить три варианта. Обычная пушка не реагирует на игрока. Motion-capture лаборатория точная, но дорогая и не для ежедневной футбольной тренировки. Project_Cam соединяет 3D-трекинг, физическую подачу мяча и контуры безопасности. Если спросят про точность: средняя ошибка мяча 95.17 mm после коррекции, средняя ошибка суставов 143.38 mm по thesis validation.

## Slide 5 — Пилот с Кайратом: что нужно
Сформулировать запрос: 3-месячный пилот с академией. Камеры — это не признание слабости, а шаг к надёжности на масштабе академии: global shutter, 60 FPS, аппаратная синхронизация, меньше motion blur, лучше быстрый мяч. Первый бюджетный ориентир: около 1.84 млн KZT на оборудование по уже подготовленному техзаданию. Завершить вопросом о пилотном формате: какие команды и возрастные группы лучше взять первыми.

## Evidence References
- Existing thesis deck converted by MarkItDown: `kairat_pitch/thesis_defense_markitdown.md`
- Project metrics: `README.md`, `CLAUDE.md`, `thesis_defense_presentation/latency_table_perf_blm_20260417_134210.png`
- Camera upgrade baseline: `docs/techspec_projectcam_ordering_ready_2026-06-04.md`
- Kairat academy scale: https://fckairat.com/academy/about/
"""
    OUT_NOTES.write_text(notes, encoding="utf-8")


def main():
    build_deck()
    build_notes()
    print(OUT_PPTX)
    print(OUT_NOTES)


if __name__ == "__main__":
    main()
